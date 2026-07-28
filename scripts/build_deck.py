"""
scripts/build_deck.py
─────────────────────
Renders the Smriti investor pitch deck (12 slides, 16:9, dark theme).

    source venv/bin/activate
    python scripts/build_deck.py

Output: smriti-deck.pptx in the project root.

Every fact on every slide is backed by the running code or by the verified
eval in scripts/eval_claims.py. No fabricated traction, quotes, customer
rosters, or market dollar figures. No em dashes.

Flow: title -> what is -> why built -> problem -> product -> why now ->
      how it works -> moat -> proof (metrics) -> demo -> traction ->
      market & model.
"""
from __future__ import annotations

import math
import os
import sys
import tempfile
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "demo_data" / "multimedia"
OUT = ROOT / "smriti-deck.pptx"

# ── Theme (dark, investor-deck calm) ─────────────────────────────────────────
BG          = RGBColor(0x0B, 0x12, 0x1F)   # near-black slate
PANEL       = RGBColor(0x14, 0x1B, 0x2D)   # card surface
PANEL_2     = RGBColor(0x1C, 0x24, 0x3A)   # raised card
ACCENT      = RGBColor(0x60, 0xA5, 0xFA)   # blue-400
ACCENT_2    = RGBColor(0xA7, 0x8B, 0xFA)   # violet-400
TEXT        = RGBColor(0xE5, 0xE7, 0xEB)   # slate-200
MUTED       = RGBColor(0x94, 0xA3, 0xB8)   # slate-400
SUCCESS     = RGBColor(0x34, 0xD3, 0x99)   # emerald-400
DANGER      = RGBColor(0xF8, 0x71, 0x71)   # red-400
GOLD        = RGBColor(0xFB, 0xBF, 0x24)   # amber-400
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SUBTLE_BD   = RGBColor(0x2A, 0x34, 0x4B)   # hairline

FONT = "Helvetica Neue"


# ── Helpers ──────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill, line=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        s.adjustments[0] = 0.06
    except Exception:
        pass
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w is not None:
            s.line.width = line_w
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *,
             size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font=FONT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=16, color=TEXT,
                bullet_color=ACCENT, line_spacing=1.25, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        r1 = p.add_run()
        r1.text = "▸  "
        r1.font.name = font
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet_color
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = item
        r2.font.name = font
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb


def add_image(slide, path, x, y, w=None, h=None):
    if w and h:
        return slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    if w:
        return slide.shapes.add_picture(str(path), x, y, width=w)
    if h:
        return slide.shapes.add_picture(str(path), x, y, height=h)
    return slide.shapes.add_picture(str(path), x, y)


def page_chrome(slide, page_no: int, total: int, title: str | None = None) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.08))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.shadow.inherit = False
    add_text(slide, Inches(0.5), Inches(0.18), Inches(6), Inches(0.35),
             "SMRITI", size=12, bold=True, color=ACCENT, font=FONT)
    add_text(slide, Inches(1.25), Inches(0.18), Inches(6), Inches(0.35),
             "internal knowledge agent", size=10, color=MUTED, font=FONT)
    add_text(slide, Inches(11.5), Inches(0.18), Inches(1.5), Inches(0.35),
             f"{page_no:02d} / {total:02d}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)
    if title:
        add_text(slide, Inches(0.6), Inches(0.55), Inches(12), Inches(0.7),
                 title, size=32, bold=True, color=WHITE)
        u = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.22),
                                   Inches(0.6), Inches(0.05))
        u.line.fill.background()
        u.fill.solid()
        u.fill.fore_color.rgb = ACCENT
        u.shadow.inherit = False


# ── Diagram renderers ────────────────────────────────────────────────────────

def render_architecture_png(path: Path) -> Path:
    """Real pipeline: sources -> parser/vision/whisper -> pgvector ->
    hybrid retrieval + cross-encoder rerank -> ReAct agent (qwen2.5:7b,
    5 read-only tools) -> two-layer grounding firewall -> cited answer or refusal."""
    W, H = 1600, 900
    img = Image.new("RGB", (W, H), (11, 18, 31))
    d = ImageDraw.Draw(img)

    F_TITLE = ImageFont.truetype("/System/Library/Fonts/HelveticaNeue.ttc", size=32)
    F_BOX   = ImageFont.truetype("/System/Library/Fonts/HelveticaNeue.ttc", size=24)
    F_SMALL = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", size=18)
    F_TINY  = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", size=15)

    def panel(x, y, w, h, title, lines, accent=(96, 165, 250)):
        d.rounded_rectangle([x, y, x + w, y + h], radius=18,
                            outline=(42, 52, 75), width=2, fill=(20, 27, 45))
        d.rounded_rectangle([x, y, x + 6, y + h], radius=4, fill=accent)
        d.text((x + 22, y + 14), title, font=F_TITLE, fill=(229, 231, 235))
        for i, line in enumerate(lines):
            color = (148, 163, 184) if line.startswith("·") else (200, 210, 230)
            d.text((x + 22, y + 58 + i * 28), line, font=F_SMALL, fill=color)

    def arrow(x1, y1, x2, y2, label=None):
        d.line([(x1, y1), (x2, y2)], fill=(96, 165, 250), width=3)
        ang = math.atan2(y2 - y1, x2 - x1)
        ah = 12
        ax, ay = x2, y2
        for off in (ang + math.pi - 0.4, ang + math.pi + 0.4):
            d.line([(ax, ay), (ax + ah * math.cos(off), ay + ah * math.sin(off))],
                   fill=(96, 165, 250), width=3)
        if label:
            mx, my = (x1 + x2) // 2, (y1 + y2) // 2
            d.text((mx + 6, my - 18), label, font=F_TINY, fill=(148, 163, 184))

    panel(60, 80, 300, 740, "Sources", [
        "Markdown",  "· memos, SOPs",
        "PDF",       "· text + scanned",
        "Images",    "· PNG / JPG / WebP",
        "Video",     "· MP4 (audio -> whisper)",
        "Audio",     "· WAV / MP3",
        "Code",      "· Python / JS / TS",
    ], accent=(96, 165, 250))

    panel(400, 80, 380, 740, "Ingestion", [
        "parser.py",
        "· chunk 1000 / 200",
        "· PII + secret scrub",
        "",
        "vision_processor.py",
        "· moondream + OCR",
        "· sidecar priority",
        "",
        "transcription.py",
        "· whisper-tiny",
        "",
        "pgvector (768-dim)",
        "· HNSW + row security",
    ], accent=(167, 139, 250))

    panel(820, 80, 400, 360, "Retrieval", [
        "hybrid search",
        "· semantic top-60",
        "· keyword top-60",
        "· RRF fusion (k=60)",
        "· cross-encoder rerank",
        "· ms-marco-MiniLM",
        "· named-doc boost",
        "· top-8 to agent",
    ], accent=(52, 211, 153))

    panel(820, 460, 400, 360, "ReAct agent (qwen2.5:7b)", [
        "loop: think / tool / observe",
        "max 5 iterations, temp 0.0",
        "",
        "Tools (read-only):",
        "· search_documents",
        "· read_chunk",
        "· list_files",
        "· compare_sections",
        "· summarize_document",
    ], accent=(248, 113, 113))

    panel(1260, 80, 300, 740, "Grounding firewall", [
        "two-layer check",
        "· verify each sentence",
        "· year + number anchor",
        "· strip unsupported",
        "· auto-cite verified",
        "· refuse if all stripped",
        "",
        "Output",
        "· answer + citations",
        "· strict refusal",
    ], accent=(250, 204, 21))

    arrow(360, 450, 400, 450)                      # Sources -> Ingestion
    arrow(780, 230, 820, 230)                       # Ingestion -> Retrieval
    arrow(1220, 230, 1260, 230, "chunks")           # Retrieval -> Grounding (evidence)
    arrow(1220, 640, 1260, 640, "answer")           # Agent -> Grounding (draft)
    arrow(1020, 440, 1020, 460, "retrieve")         # Retrieval -> Agent
    # iterate arc: agent back up to retrieval
    d.arc([760, 420, 830, 500], start=90, end=270, fill=(248, 113, 113), width=2)
    d.text((720, 432), "iterate", font=F_TINY, fill=(248, 113, 113))

    img.save(path, optimize=True)
    return path


def render_market_png(path: Path) -> Path:
    """Buyer segments and the wedge. No dollar figures (no fabricated market sizing)."""
    W, H = 1400, 900
    img = Image.new("RGB", (W, H), (11, 18, 31))
    d = ImageDraw.Draw(img)

    F_BIG   = ImageFont.truetype("/System/Library/Fonts/HelveticaNeue.ttc", size=40)
    F_MED   = ImageFont.truetype("/System/Library/Fonts/HelveticaNeue.ttc", size=30)
    F_LABEL = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", size=26)
    F_NOTE  = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", size=20)
    F_TINY  = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", size=18)

    d.text((40, 40), "Who buys", font=F_BIG, fill=(229, 231, 235))
    d.text((40, 92), "Regulated enterprises that cannot ship a hallucinating AI.",
           font=F_NOTE, fill=(148, 163, 184))

    def card(x, y, w, h, title, body, accent):
        d.rounded_rectangle([x, y, x + w, y + h], radius=16,
                            outline=(42, 52, 75), width=2, fill=(20, 27, 45))
        d.rounded_rectangle([x, y, x + 6, y + h], radius=4, fill=accent)
        d.text((x + 22, y + 18), title, font=F_MED, fill=(229, 231, 235))
        for i, line in enumerate(body):
            d.text((x + 22, y + 64 + i * 30), line, font=F_TINY, fill=(200, 210, 230))

    seg = [
        (40, 160, 640, 200, "Banks & NBFCs",
         ["RBI / SEBI supervised.", "Audit trails and refusal logs expected."], (96, 165, 250)),
        (700, 160, 660, 200, "Insurers",
         ["IRDAI model governance.", "Human-in-loop for material decisions."], (167, 139, 250)),
        (40, 400, 640, 200, "Medtech (HIPAA)",
         ["On-prem mandatory.", "No PHI egress, no sub-processors."], (52, 211, 153)),
        (700, 400, 660, 200, "BPOs for regulated clients",
         ["Serve US banks and insurers.", "Need citeable, auditable answers."], (248, 113, 113)),
    ]
    for s in seg:
        card(*s)

    # Wedge banner
    wx, wy, ww, wh = 40, 640, 1320, 170
    d.rounded_rectangle([wx, wy, wx + ww, wy + wh], radius=16,
                        outline=(250, 204, 21), width=2, fill=(26, 26, 16))
    d.rounded_rectangle([wx, wy, wx + 6, wy + wh], radius=4, fill=(250, 204, 21))
    d.text((wx + 22, wy + 16), "The wedge", font=F_MED, fill=(250, 204, 21))
    d.text((wx + 22, wy + 60), "Compliance and audit need forces three things: citation, refusal, and on-prem.",
           font=F_NOTE, fill=(229, 231, 235))
    d.text((wx + 22, wy + 96), "Generic copilots fail all three. Smriti is built for exactly this.",
           font=F_NOTE, fill=(200, 210, 230))

    img.save(path, optimize=True)
    return path


# ── Slide builders ───────────────────────────────────────────────────────────

def slide_title(prs, total: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG)

    add_text(s, Inches(0.8), Inches(1.5), Inches(12), Inches(1.5),
             "Smriti", size=120, bold=True, color=WHITE, font=FONT)
    add_text(s, Inches(0.85), Inches(3.1), Inches(12), Inches(0.7),
             "The internal knowledge agent that refuses to guess.",
             size=32, color=ACCENT, font=FONT)
    add_text(s, Inches(0.85), Inches(3.85), Inches(12), Inches(0.5),
             "Multimodal RAG for banks, insurers, and other regulated enterprises.",
             size=20, color=MUTED, italic=True, font=FONT)

    add_rect(s, Inches(0.85), Inches(5.0), Inches(11.6), Inches(1.3), PANEL, line=SUBTLE_BD, line_w=Inches(0.01))
    add_text(s, Inches(1.05), Inches(5.15), Inches(11), Inches(0.4),
             "Why now", size=12, color=MUTED, font=FONT)
    add_text(s, Inches(1.05), Inches(5.5), Inches(11), Inches(0.7),
             "Compliance teams can't ship copilots that hallucinate.   We built one that doesn't.",
             size=22, bold=True, color=WHITE, font=FONT)

    add_text(s, Inches(0.85), Inches(6.7), Inches(12), Inches(0.4),
             "2026", size=14, color=MUTED, italic=True, font=FONT)
    page_chrome(s, 1, total)


def slide_what_is(prs, total: int, n: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG)
    page_chrome(s, n, total, "What is Smriti")

    add_text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
             "Plain English, no acronyms.",
             size=20, color=MUTED, italic=True)

    add_rect(s, Inches(0.6), Inches(2.2), Inches(12.13), Inches(4.3), PANEL,
             line=SUBTLE_BD, line_w=Inches(0.01))
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.2), Inches(0.1), Inches(4.3))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.shadow.inherit = False

    paras = [
        ("An AI search box",
         "for everything your company has ever written, screenshotted, or recorded."),
        ("One answer, one citation",
         "an employee types a question in plain English and gets a single answer, with a citation to the exact file and section it came from."),
        ("It says no when it should",
         "if the answer isn't in the company's own documents, Smriti says so. It does not make one up."),
        ("It runs on the customer's own servers",
         "in their VPC, behind their firewall. No data ever leaves. No SaaS, no third-party API."),
    ]
    y = Inches(2.5)
    for headline, body in paras:
        add_text(s, Inches(1.0), y, Inches(11.5), Inches(0.4),
                 headline, size=18, bold=True, color=ACCENT, font=FONT)
        add_text(s, Inches(1.0), y + Inches(0.42), Inches(11.5), Inches(0.45),
                 body, size=15, color=TEXT, font=FONT)
        y += Inches(0.95)

    add_text(s, Inches(0.6), Inches(6.7), Inches(12.13), Inches(0.4),
             "Not ChatGPT for work. Not a Copilot wrapper. Not a vector DB with a UI. "
             "A grounded, auditable, on-prem knowledge agent.",
             size=13, color=MUTED, italic=True, align=PP_ALIGN.CENTER, font=FONT)


def slide_why_built(prs, total: int, n: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG)
    page_chrome(s, n, total, "Why we built it")

    add_text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
             "I was the analyst waiting days for an answer. Smriti is what I wished existed.",
             size=20, color=MUTED, italic=True)

    add_rect(s, Inches(0.6), Inches(2.2), Inches(6.4), Inches(4.9), PANEL,
             line=SUBTLE_BD, line_w=Inches(0.01))
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.2), Inches(0.1), Inches(4.9))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_2
    bar.shadow.inherit = False

    add_text(s, Inches(0.85), Inches(2.35), Inches(6), Inches(0.4),
             "The founder's story", size=12, color=MUTED, font=FONT)

    paras = [
        "During Covid I worked remotely at Accenture on insurance claims workflows. "
        "I had no idea how those workflows actually ran. The knowledge lived in "
        "senior heads, not anywhere I could reach.",
        "I would wait whole days for a senior to answer a single question. They "
        "wanted to help, but they were tangled in their own deliverables and could not.",
        "Worse, I kept hitting the same errors other developers had already hit. "
        "Each one needed a senior to confirm the fix. More days lost, on repeat.",
        "Smriti is what I wished existed then: the team's knowledge in one place, "
        "answerable in seconds, with sensitive data never leaving the company intranet.",
    ]
    y = Inches(2.8)
    for body in paras:
        add_text(s, Inches(0.95), y, Inches(5.95), Inches(1.05),
                 body, size=13, color=TEXT, font=FONT)
        y += Inches(1.02)

    add_rect(s, Inches(7.25), Inches(2.2), Inches(5.48), Inches(4.9), PANEL_2,
             line=SUBTLE_BD, line_w=Inches(0.01))
    add_text(s, Inches(7.5), Inches(2.35), Inches(5), Inches(0.4),
             "What this means", size=12, color=MUTED, font=FONT)
    add_bullets(s, Inches(7.5), Inches(2.95), Inches(5.1), Inches(4.0),
                [
                    "The pain is personal, not researched. The founder was the blocked analyst.",
                    "Built for the person waiting on a senior, not the exec buying a tool.",
                    "Data sovereignty is the founding constraint, not a feature bolted on later.",
                    "AI is now widely available. The open question is who owns how your data is used. Smriti keeps that ownership in the customer's hands.",
                ],
                size=13, color=TEXT, bullet_color=GOLD, line_spacing=1.3)


def slide_problem(prs, total: int, n: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG)
    page_chrome(s, n, total, "The problem")

    add_text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
             "Every regulated enterprise has the same problem.",
             size=20, color=MUTED, italic=True)

    cards = [
        ("Knowledge is everywhere",
         "RBI memos in drives. KYC flow in Confluence. Org chart in a PNG. "
         "Compliance briefing in a video. None of it is searchable together.",
         DANGER),
        ("Employees can't find it",
         "Finding the right document is slow and manual, scattered across "
         "formats and systems. A regulated answer needs a citation, not a vibe.",
         ACCENT),
        ("Generic copilots can't ship",
         "Off-the-shelf LLMs hallucinate regulatory citations. "
         "Compliance teams block the rollout. AI budgets stay unspent.",
         ACCENT_2),
    ]
    for i, (title, body, accent) in enumerate(cards):
        x = Inches(0.6 + i * 4.18)
        add_rect(s, x, Inches(2.4), Inches(3.95), Inches(4.4), PANEL, line=SUBTLE_BD, line_w=Inches(0.01))
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.4), Inches(0.1), Inches(4.4))
        bar.line.fill.background()
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.shadow.inherit = False
        add_text(s, x + Inches(0.3), Inches(2.55), Inches(3.5), Inches(0.6),
                 title, size=20, bold=True, color=WHITE)
        add_text(s, x + Inches(0.3), Inches(3.2), Inches(3.5), Inches(3.4),
                 body, size=14, color=TEXT)


def slide_solution(prs, total: int, n: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG)
    page_chrome(s, n, total, "The product")

    add_text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
             "One agent that reads everything the company knows, and only answers from it.",
             size=20, color=MUTED, italic=True)

    add_rect(s, Inches(0.6), Inches(2.4), Inches(5.4), Inches(4.9), PANEL, line=SUBTLE_BD, line_w=Inches(0.01))
    add_text(s, Inches(0.85), Inches(2.6), Inches(5), Inches(0.5),
             "What it does", size=22, bold=True, color=WHITE)
    add_bullets(s, Inches(0.85), Inches(3.15), Inches(5), Inches(4),
                [
                    "Indexes markdown, PDFs, images, video, audio, and code into one pgvector store",
                    "Hybrid retrieval: dense vectors plus full-text, fused and re-ranked",
                    "ReAct agent reasons across the corpus with five read-only tools",
                    "Two-layer grounding check on every answer, claims must appear in source",
                    "Refuses cleanly when the answer isn't in the corpus",
                    "Returns a citation for every fact, not a vibe",
                ],
                size=14, color=TEXT, bullet_color=ACCENT, line_spacing=1.2)

    add_text(s, Inches(6.4), Inches(2.4), Inches(6.4), Inches(0.4),
             "Reads screenshots, PDFs, and narrated video, like a person would.",
             size=14, color=MUTED, italic=True)
    add_image(s, ASSETS / "dashboard-mrr-q2.png",
              Inches(6.4), Inches(2.85), w=Inches(6.4))


def slide_why_now(prs, total: int, n: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG)
    page_chrome(s, n, total, "Why now")

    add_text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
             "The window is open, on both sides of the ocean.",
             size=20, color=MUTED, italic=True)

    cols = [
        ("United States",
         "Federal + state + sectoral",
         [
             "OCC SR 11-7 (model risk management): examiners now apply it to third-party AI. A bank that ships a copilot with no audit trail or refusal log risks a Matter Requiring Attention.",
             "SEC cybersecurity disclosure rules: put the CISO on the hook for AI data flows. Copilots that call home fail on day one. Smriti is single-tenant, no egress.",
             "NIST AI RMF plus the GenAI Profile: the de facto standard for federal procurement and the largest US banks. Smriti maps directly to VALIDATE and VERIFY.",
             "State insurance regulators (NAIC model bulletin on AI): pushing insurers toward documented AI governance and impact assessments.",
         ],
         ACCENT),
        ("India",
         "RBI + SEBI + IRDAI",
         [
             "RBI: digital lending and AI/ML supervisory expectations now require explainability, auditability, and an AI use-case inventory for NBFCs.",
             "RBI expects a documented review process for any AI-assisted decisioning on credit, KYC, or customer service. A hallucinated answer fails this on its face.",
             "SEBI: AI/ML in markets is pushing broker-dealers, mutual funds, and RTAs to formalize AI governance. Smriti's citation-first, refusal-aware design fits directly.",
             "IRDAI: AI use by insurers requires documented model governance, periodic validation, and a human-in-the-loop for material decisions.",
         ],
         ACCENT_2),
    ]
    for i, (headline, badge, bullets, accent) in enumerate(cols):
        x = Inches(0.6 + i * 6.12)
        add_rect(s, x, Inches(2.4), Inches(5.9), Inches(4.7), PANEL, line=SUBTLE_BD, line_w=Inches(0.01))
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.4), Inches(0.1), Inches(4.7))
        bar.line.fill.background()
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.shadow.inherit = False
        add_rect(s, x + Inches(0.3), Inches(2.55), Inches(2.5), Inches(0.45), accent)
        add_text(s, x + Inches(0.3), Inches(2.6), Inches(2.5), Inches(0.35),
                 badge, size=11, bold=True, color=BG, align=PP_ALIGN.CENTER, font=FONT)
        add_text(s, x + Inches(0.3), Inches(3.2), Inches(5.4), Inches(0.6),
                 headline, size=22, bold=True, color=WHITE)
        add_bullets(s, x + Inches(0.3), Inches(3.95), Inches(5.4), Inches(3.05),
                    bullets,
                    size=11, color=TEXT, bullet_color=accent, line_spacing=1.15)


def slide_how(prs, total: int, n: int, arch_path: Path) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG)
    page_chrome(s, n, total, "How it works")

    add_text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
             "Sources to ingest to vector store to ReAct agent to grounded response.",
             size=18, color=MUTED, italic=True)

    add_image(s, arch_path, Inches(2.22), Inches(2.15), w=Inches(8.89), h=Inches(5.0))


def slide_moat(prs, total: int, n: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG)
    page_chrome(s, n, total, "The moat")

    add_text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
             "Anyone can wrap an LLM. We make it safe to put in front of a regulator.",
             size=20, color=MUTED, italic=True)

    add_rect(s, Inches(0.6), Inches(2.4), Inches(6.0), Inches(4.9), PANEL, line=SUBTLE_BD, line_w=Inches(0.01))
    add_text(s, Inches(0.85), Inches(2.6), Inches(5.5), Inches(0.5),
             "Four things that don't exist in a wrapper", size=20, bold=True, color=WHITE)
    add_bullets(s, Inches(0.85), Inches(3.15), Inches(5.5), Inches(4.0),
                [
                    "Claim-level grounding: every number, date, and name is verified against the source, then a second guard re-checks concrete claims",
                    "Strict refusal: no answer beats a wrong answer in a regulated workflow",
                    "Sidecar priority: human-written markdown wins over the vision model's interpretation",
                    "On-prem, no egress: data stays in the customer's VPC, no telemetry, no sub-processors",
                ],
                size=14, color=TEXT, bullet_color=GOLD, line_spacing=1.25)

    add_rect(s, Inches(6.85), Inches(2.4), Inches(5.95), Inches(4.9), PANEL_2, line=SUBTLE_BD, line_w=Inches(0.01))
    add_text(s, Inches(7.1), Inches(2.55), Inches(5.5), Inches(0.4),
             "What it looks like in practice", size=12, color=MUTED, italic=True)

    add_text(s, Inches(7.1), Inches(3.0), Inches(5.5), Inches(0.4),
             "Q: What was our Q2 MRR figure in crores?",
             size=14, color=ACCENT, font="Menlo")
    add_text(s, Inches(7.1), Inches(3.5), Inches(5.5), Inches(1.5),
             "“Q2 MRR was INR 4.8 Crore, up 22% QoQ. "
             "[Source: multimedia/dashboard-mrr-q2.md | Section 1]”",
             size=13, color=WHITE, italic=True, font=FONT)

    add_text(s, Inches(7.1), Inches(4.6), Inches(5.5), Inches(0.4),
             "Q: What is the revenue of Tesla?",
             size=14, color=ACCENT, font="Menlo")
    add_text(s, Inches(7.1), Inches(5.1), Inches(5.5), Inches(1.5),
             "“I don't have that information from the indexed documents. "
             "No relevant content was found that supports an answer.”",
             size=13, color=DANGER, italic=True, font=FONT)


def slide_metrics(prs, total: int, n: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG)
    page_chrome(s, n, total, "The proof: measured, not claimed")

    add_text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
             "64-question evaluation on a public 819-document benchmark. Zero fabricated answers.",
             size=20, color=MUTED, italic=True)

    add_rect(s, Inches(0.6), Inches(2.05), Inches(12.13), Inches(1.45), PANEL_2,
             line=SUBTLE_BD, line_w=Inches(0.01))
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.05), Inches(0.1), Inches(1.45))
    bar.line.fill.background(); bar.fill.solid(); bar.fill.fore_color.rgb = SUCCESS; bar.shadow.inherit = False
    add_text(s, Inches(0.95), Inches(2.15), Inches(3.5), Inches(1.3),
             "100%", size=64, bold=True, color=SUCCESS, font=FONT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(4.6), Inches(2.25), Inches(8), Inches(0.5),
             "Never made things up.", size=24, bold=True, color=WHITE, font=FONT)
    add_text(s, Inches(4.6), Inches(2.75), Inches(8), Inches(0.7),
             "Across 64 questions and 4 test types, every answer was either cited to a real "
             "source or honestly refused. Zero invented facts, the property regulated buyers need.",
             size=13, color=TEXT, font=FONT)

    cards = [
        ("Answerable", "30 Qs", "100%", "Cited the real source, or said “I don't have that.”"),
        ("Out-of-corpus", "19 Qs", "100%", "Refused to guess on things it couldn't know."),
        ("Trick questions", "12 Qs", "100%", "Real topic, absent detail, no invented specifics."),
        ("Two-document", "3 Qs", "100%", "Synthesized and cited both required sources."),
    ]
    for i, (title, count, pct, body) in enumerate(cards):
        x = Inches(0.6 + i * 3.04)
        add_rect(s, x, Inches(3.7), Inches(2.95), Inches(1.7), PANEL, line=SUBTLE_BD, line_w=Inches(0.01))
        add_text(s, x + Inches(0.22), Inches(3.82), Inches(2.6), Inches(0.35),
                 title, size=13, bold=True, color=WHITE, font=FONT)
        add_text(s, x + Inches(0.22), Inches(4.12), Inches(2.6), Inches(0.3),
                 count, size=11, color=MUTED, font=FONT)
        add_text(s, x + Inches(0.22), Inches(4.42), Inches(2.6), Inches(0.55),
                 pct, size=30, bold=True, color=SUCCESS, font=FONT)
        add_text(s, x + Inches(0.22), Inches(4.95), Inches(2.6), Inches(0.4),
                 body, size=10, color=TEXT, font=FONT)

    add_rect(s, Inches(0.6), Inches(5.6), Inches(12.13), Inches(1.05), PANEL, line=SUBTLE_BD, line_w=Inches(0.01))
    add_text(s, Inches(0.85), Inches(5.72), Inches(7), Inches(0.5),
             "Found the right document: 90%", size=20, bold=True, color=ACCENT, font=FONT)
    add_text(s, Inches(0.85), Inches(6.12), Inches(8.5), Inches(0.5),
             "The 10% are wrong-source citations, real documents, just not the perfect one. "
             "Never a fake source. Improving: 53% then 73% then 90%.",
             size=11, color=TEXT, font=FONT)
    add_text(s, Inches(9.0), Inches(5.78), Inches(3.6), Inches(0.8),
             "A weaker Smriti refuses more.\nIt never degrades into guessing.",
             size=12, bold=True, color=GOLD, align=PP_ALIGN.RIGHT, font=FONT, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.6), Inches(6.85), Inches(12.13), Inches(0.4),
             "Public benchmark, not customer data; per-pilot numbers established on each customer's corpus. "
             "Recall not yet formally measured; answer-coverage proxy reported instead.",
             size=10, color=MUTED, italic=True, align=PP_ALIGN.CENTER, font=FONT)


def slide_demo(prs, total: int, n: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG)
    page_chrome(s, n, total, "Live demo")

    add_text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
             "A real query, end to end, fully grounded.",
             size=20, color=MUTED, italic=True)

    steps = [
        ("1. Ask", "“What is the current status of the Pine Labs webhook incident?”",
         "Routes to the agent over a single REST call.", ACCENT),
        ("2. Retrieve", "search_documents(query='Pine Labs webhook', category='multimedia')\n"
                          "→ top hit: incident-status-board.md\n"
                          "→ read_chunk → verbatim text",
         "Hybrid retrieval, not just a vector lookup.", ACCENT_2),
        ("3. Answer", "“The Pine Labs webhook outage is P1 OPEN with an ETA of "
                       "approximately 4 hours from July 15, 2025. "
                       "[Source: multimedia/incident-status-board.md | Section 1]”",
         "Grounding check passes, every claim is verbatim in the source.", SUCCESS),
    ]
    for i, (label, body, sub, accent) in enumerate(steps):
        y = Inches(2.4 + i * 1.55)
        add_rect(s, Inches(0.6), y, Inches(12.13), Inches(1.4), PANEL, line=SUBTLE_BD, line_w=Inches(0.01))
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(0.1), Inches(1.4))
        bar.line.fill.background()
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.shadow.inherit = False
        add_text(s, Inches(0.85), y + Inches(0.1), Inches(2.2), Inches(0.4),
                 label, size=14, bold=True, color=accent, font=FONT)
        add_text(s, Inches(3.0), y + Inches(0.12), Inches(9.5), Inches(0.8),
                 body, size=12, color=TEXT, font="Menlo")
        add_text(s, Inches(0.85), y + Inches(1.0), Inches(11.7), Inches(0.35),
                 sub, size=11, color=MUTED, italic=True)


def slide_traction(prs, total: int, n: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG)
    page_chrome(s, n, total, "Traction & next steps")

    add_text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
             "Working product, one warm enterprise lead, compliance-ready architecture.",
             size=20, color=MUTED, italic=True)

    add_rect(s, Inches(0.6), Inches(2.4), Inches(3.95), Inches(1.7), PANEL, line=SUBTLE_BD, line_w=Inches(0.01))
    add_text(s, Inches(0.85), Inches(2.55), Inches(3.5), Inches(0.4),
             "Product", size=12, color=MUTED, font=FONT)
    add_text(s, Inches(0.85), Inches(2.85), Inches(3.5), Inches(1.0),
             "Live", size=48, bold=True, color=SUCCESS, font=FONT)
    add_text(s, Inches(0.85), Inches(3.65), Inches(3.5), Inches(0.4),
             "Multimodal RAG, on-prem, end to end.", size=12, color=TEXT, font=FONT)

    add_rect(s, Inches(4.7), Inches(2.4), Inches(3.95), Inches(1.7), PANEL, line=SUBTLE_BD, line_w=Inches(0.01))
    add_text(s, Inches(4.95), Inches(2.55), Inches(3.5), Inches(0.4),
             "Compliance-ready", size=12, color=MUTED, font=FONT)
    add_text(s, Inches(4.95), Inches(2.85), Inches(3.5), Inches(1.0),
             "Built", size=48, bold=True, color=ACCENT, font=FONT)
    add_text(s, Inches(4.95), Inches(3.65), Inches(3.5), Inches(0.4),
             "OIDC SSO, RBAC, audit log, SBOM, no egress.", size=12, color=TEXT, font=FONT)

    add_rect(s, Inches(8.8), Inches(2.4), Inches(3.95), Inches(1.7), PANEL, line=SUBTLE_BD, line_w=Inches(0.01))
    add_text(s, Inches(9.05), Inches(2.55), Inches(3.5), Inches(0.4),
             "Enterprise lead", size=12, color=MUTED, font=FONT)
    add_text(s, Inches(9.05), Inches(2.85), Inches(3.5), Inches(1.0),
             "1", size=48, bold=True, color=ACCENT_2, font=FONT)
    add_text(s, Inches(9.05), Inches(3.65), Inches(3.5), Inches(0.4),
             "Global medtech (CGM), HIPAA-regulated.", size=12, color=TEXT, font=FONT)

    add_rect(s, Inches(0.6), Inches(4.4), Inches(12.15), Inches(2.8), PANEL, line=SUBTLE_BD, line_w=Inches(0.01))
    add_text(s, Inches(0.85), Inches(4.55), Inches(11.5), Inches(0.5),
             "Where we are", size=18, bold=True, color=WHITE)
    add_bullets(s, Inches(0.95), Inches(5.1), Inches(11.5), Inches(2.0),
                [
                    "Working product: 819-doc public benchmark, 64-question eval, 100% cite-or-refuse, 90% retrieval precision.",
                    "One warm enterprise lead: a global medtech / CGM company, HIPAA-regulated, wants a compliance review before any real data.",
                    "Pilot scope proposed: internal knowledge retrieval over non-clinical documents (SOPs, quality, policies), not clinical decision support.",
                    "Compliance gaps already closed: OIDC/JWKS SSO, admin/user RBAC, append-only audit log, SBOM, security posture doc.",
                ],
                size=12, color=TEXT, bullet_color=SUCCESS, line_spacing=1.2)


def slide_market_model(prs, total: int, n: int, mkt_path: Path) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG)
    page_chrome(s, n, total, "Market & business model")

    add_text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
             "Per-seat pricing, anchored on the audit need.",
             size=20, color=MUTED, italic=True)

    add_image(s, mkt_path, Inches(0.6), Inches(2.15), w=Inches(8.5), h=Inches(5.0))

    add_rect(s, Inches(9.4), Inches(2.4), Inches(3.4), Inches(4.7), PANEL, line=SUBTLE_BD, line_w=Inches(0.01))
    add_text(s, Inches(9.6), Inches(2.55), Inches(3.1), Inches(0.4),
             "Proposed pricing", size=12, color=MUTED, font=FONT)
    add_text(s, Inches(9.6), Inches(2.85), Inches(3.1), Inches(0.7),
             "$30 / seat / mo", size=24, bold=True, color=ACCENT, font=FONT)
    add_text(s, Inches(9.6), Inches(3.45), Inches(3.1), Inches(0.4),
             "Minimum $60K annual contract.", size=11, color=MUTED, italic=True, font=FONT)

    add_text(s, Inches(9.6), Inches(4.0), Inches(3.1), Inches(0.4),
             "Land & expand", size=12, color=MUTED, font=FONT)
    add_bullets(s, Inches(9.6), Inches(4.35), Inches(3.1), Inches(2.6),
                [
                    "Land: 1 team, ~50 seats",
                    "Expand: compliance, legal, ops",
                    "Upsell: on-prem hosting + audit API",
                ],
                size=12, color=TEXT, bullet_color=SUCCESS, line_spacing=1.2)


# ── Main ─────────────────────────────────────────────────────────────────────

def main() -> int:
    with tempfile.TemporaryDirectory() as td:
        arch_path = Path(td) / "architecture.png"
        mkt_path  = Path(td) / "market.png"
        render_architecture_png(arch_path)
        render_market_png(mkt_path)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        TOTAL = 12
        slide_title(prs, TOTAL)                # 1
        slide_what_is(prs, TOTAL, 2)           # 2
        slide_why_built(prs, TOTAL, 3)         # 3
        slide_problem(prs, TOTAL, 4)           # 4
        slide_solution(prs, TOTAL, 5)          # 5
        slide_why_now(prs, TOTAL, 6)           # 6
        slide_how(prs, TOTAL, 7, arch_path)    # 7
        slide_moat(prs, TOTAL, 8)              # 8
        slide_metrics(prs, TOTAL, 9)           # 9
        slide_demo(prs, TOTAL, 10)             # 10
        slide_traction(prs, TOTAL, 11)         # 11
        slide_market_model(prs, TOTAL, 12, mkt_path)  # 12

        prs.save(OUT)

    size_kb = OUT.stat().st_size // 1024
    print(f"wrote {OUT}  ({size_kb} KB)")
    return 0


if __name__ == "__main__":
    sys.exit(main())